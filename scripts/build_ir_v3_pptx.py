"""AI-MBTI IR/PR v3 — PowerPoint 발표용 (16:9).

실행: python3 scripts/build_ir_v3_pptx.py
출력: docs/IR_AI-MBTI_v3.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
PUBLIC = ROOT / "public"
OUT = ROOT / "docs" / "IR_AI-MBTI_v3.pptx"

# 컬러
NAVY = RGBColor(0x0A, 0x1B, 0x3D)
ACCENT = RGBColor(0x1E, 0x6B, 0xFF)
SOFT = RGBColor(0xF3, 0xF6, 0xFC)
BORDER = RGBColor(0xD6, 0xDE, 0xEC)
TEXT = RGBColor(0x1A, 0x1F, 0x2E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
SUCCESS = RGBColor(0x0E, 0xA3, 0x6B)
WARN = RGBColor(0xE0, 0xA8, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = WHITE

# 16:9 A4 landscape
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)

FONT = "맑은 고딕"


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_image(slide, rel, x, y, w, h=None):
    img = PUBLIC / rel
    if not img.exists():
        return None
    if h is None:
        return slide.shapes.add_picture(str(img), x, y, width=w)
    return slide.shapes.add_picture(str(img), x, y, width=w, height=h)


def add_table(slide, x, y, w, h, headers, rows, highlight=None):
    """highlight: 0-indexed row index (in data rows) to bold."""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    tbl_shp = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shp.table

    # 헤더
    for j, head in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.margin_left = Cm(0.15)
        tf.margin_right = Cm(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = head
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = FONT

    # 데이터 행
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if i % 2 == 0 else WHITE
            tf = cell.text_frame
            tf.margin_left = Cm(0.15)
            tf.margin_right = Cm(0.15)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(10)
            is_highlight = (highlight is not None and i == highlight)
            r.font.bold = is_highlight
            r.font.color.rgb = NAVY if is_highlight else TEXT
            r.font.name = FONT


def add_kpi_card(slide, x, y, w, h, label, big, sub):
    add_rect(slide, x, y, w, h, fill=SOFT, line=BORDER)
    add_text(slide, x, y + Cm(0.3), w, Cm(0.8), label,
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Cm(1.1), w, Cm(2.0), big,
             size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Cm(3.3), w, Cm(0.8), sub,
             size=10, color=SUCCESS, align=PP_ALIGN.CENTER)


def section_bar(slide, label, title, subtitle=None):
    # 상단 네이비 바
    add_rect(slide, Cm(1.5), Cm(1.0), Cm(0.8), Cm(0.3), fill=ACCENT)
    add_text(slide, Cm(1.5), Cm(1.4), Cm(20), Cm(0.8), label,
             size=11, bold=True, color=ACCENT)
    add_text(slide, Cm(1.5), Cm(2.0), Cm(30), Cm(1.8), title,
             size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Cm(1.5), Cm(4.0), Cm(30), Cm(2.0), subtitle,
                 size=12, color=TEXT)


def footer(slide, page_num):
    add_text(slide, Cm(1.5), SLIDE_H - Cm(0.9), Cm(20), Cm(0.6),
             "AI-MBTI · IR/PR v3 · mcodegc.com",
             size=9, color=MUTED)
    add_text(slide, SLIDE_W - Cm(3), SLIDE_H - Cm(0.9), Cm(2), Cm(0.6),
             str(page_num), size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════
# 슬라이드 빌더
# ═══════════════════════════════════════════════

def build():
    prs = new_prs()

    # ─── 1. 표지 ───
    s = blank(prs)
    # 로고
    add_image(s, "Adv/로고 및 참조/로고_transparent.png",
              x=(SLIDE_W - Cm(5)) / 2, y=Cm(3), w=Cm(5))
    add_text(s, Cm(0), Cm(8), SLIDE_W, Cm(1.2),
             "AI ERA CAREER DIAGNOSIS",
             size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Cm(0), Cm(9.2), SLIDE_W, Cm(3),
             "AI-MBTI",
             size=66, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Cm(0), Cm(13), SLIDE_W, Cm(1.2),
             "AI 시대, 당신의 생존력을 진단합니다",
             size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Cm(0), Cm(14.5), SLIDE_W, Cm(1),
             "20문항으로 AI 활용 유형 진단 · 데이터 기반 광고 운영",
             size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Cm(0), SLIDE_H - Cm(1.5), SLIDE_W, Cm(0.8),
             "IR / PR Deck v3 · 2026.04 · mcodegc.com",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # ─── 2. 성과 (KPI 4개) ───
    s = blank(prs)
    section_bar(s, "01  성과", "숫자로 보는 현재",
                "2026년 4월 22일 ~ 4월 23일, 페이스북·인스타 광고 3종 집행 결과")
    kpi_y = Cm(7)
    kpi_h = Cm(4.2)
    kpi_w = Cm(7)
    kpi_gap = Cm(0.5)
    start_x = (SLIDE_W - (kpi_w * 4 + kpi_gap * 3)) / 2
    for i, (lbl, big, sub) in enumerate([
        ("광고비", "₩88,813", "2일간 집행"),
        ("관심 고객", "17명", "중복 제거"),
        ("1명 모시는 비용", "₩4,030", "최우수 광고"),
        ("광고 클릭률", "3.03%", "업계의 1.6배"),
    ]):
        add_kpi_card(s, start_x + (kpi_w + kpi_gap) * i, kpi_y, kpi_w, kpi_h, lbl, big, sub)
    footer(s, 2)

    # ─── 3. 캠페인 3종 성적표 ───
    s = blank(prs)
    add_text(s, Cm(1.5), Cm(1), Cm(30), Cm(1.2),
             "캠페인 3종 성적표 (여성 이미지 × 첫페이지 3종 · 4월 22~23일)",
             size=20, bold=True, color=NAVY)
    add_table(s, Cm(1.5), Cm(2.8), Cm(30.5), Cm(6),
              ["캠페인 (여성 × 첫페이지)", "관심 고객", "1명 모시는 비용", "한 줄 평가"],
              [
                  ["여성 × 심플 첫페이지 (1등)", "7명", "₩4,030", "정보 중심 · 가장 저렴"],
                  ["여성 × 사회적 첫페이지", "7명", "₩4,121", "공감형 · 클릭 가장 많음"],
                  ["여성 × 공포 첫페이지 (실험)", "3명", "₩10,585", "효과 애매, 중단 예정"],
                  ["합계 (3종)", "17명", "평균 ₩5,224", "심플·사회적 중심 운영"],
              ],
              highlight=0)

    add_text(s, Cm(1.5), Cm(10), Cm(30), Cm(1),
             "우리 광고 vs 업계 평균",
             size=18, bold=True, color=NAVY)
    add_table(s, Cm(1.5), Cm(11.5), Cm(30.5), Cm(5),
              ["지표", "업계 평균", "우리", "차이"],
              [
                  ["광고 보고 누르는 비율", "1.91%", "3.03%", "▲ 1.6배"],
                  ["한 번 누를 때 광고비", "₩1,100", "₩562", "▼ 절반"],
                  ["1명 모시는 비용", "약 ₩25,000", "₩4,030", "▼ 1/6"],
              ],
              highlight=2)
    footer(s, 3)

    # ─── 4. 관심 고객 + 규모별 광고비 ───
    s = blank(prs)
    add_text(s, Cm(1.5), Cm(1), Cm(30), Cm(1.2),
             "관심 고객 17명 — 어디로 갔나? (4월 22~23일)",
             size=20, bold=True, color=NAVY)
    add_table(s, Cm(1.5), Cm(2.8), Cm(30.5), Cm(5),
              ["캠페인", "총 관심 고객", "전자책 이동", "단톡방 이동"],
              [
                  ["Simple", "7명", "7명", "0명"],
                  ["Social", "7명", "6명", "1명"],
                  ["Fear", "3명", "2명", "1명"],
                  ["합계", "17명", "15명", "2명"],
              ],
              highlight=3)

    add_text(s, Cm(1.5), Cm(9), Cm(30), Cm(1),
             "규모별 광고비 환산 — 고객 N명 모으려면 얼마?",
             size=18, bold=True, color=NAVY)
    add_table(s, Cm(1.5), Cm(10.5), Cm(30.5), Cm(6),
              ["규모", "현재 CPL 기준 (₩4,030)", "실질 가격 (150%)", "해석"],
              [
                  ["10명", "₩40,300", "₩60,450", "테스트"],
                  ["100명", "₩403,000", "₩604,500", "월 소규모"],
                  ["1,000명", "₩4,030,000", "₩6,045,000", "월 중형"],
                  ["10,000명", "₩40,300,000", "₩60,450,000", "연간 대규모"],
              ],
              highlight=3)
    footer(s, 4)

    # ─── 5. Meta 퍼널 비교 + 결론 ───
    s = blank(prs)
    add_text(s, Cm(1.5), Cm(1), Cm(30), Cm(1.2),
             "Meta 광고 캠페인별 전체 퍼널 비교 (4월 22~23일)",
             size=20, bold=True, color=NAVY)
    add_table(s, Cm(1.5), Cm(2.8), Cm(30.5), Cm(4.5),
              ["캠페인", "랜딩 조회", "CTA 클릭", "테스트 시작", "결과 확인", "Lead",
               "CTA 클릭률 ★", "Lead 전환율 ★"],
              [
                  ["Simple", "92", "51", "44", "37", "7", "55.4% 🏆 1위", "7.6%"],
                  ["Social", "73", "36", "33", "32", "7", "49.3%", "9.6% 🏆 1위"],
                  ["Fear", "84", "45", "38", "29", "3", "53.6%", "3.6% (꼴찌)"],
                  ["합계·평균", "249", "132", "115", "98", "17", "53.0%", "6.8%"],
              ],
              highlight=3)
    add_text(s, Cm(1.5), Cm(7.5), Cm(30.5), Cm(1),
             "★ CTA 클릭률 1위 = Simple (55.4%) / Lead 전환율 1위 = Social (9.6%) — 서로 다른 1위",
             size=10, color=MUTED)

    # 결론 3박스
    box_y = Cm(9.5)
    box_w = Cm(10)
    box_h = Cm(2.5)
    box_gap = Cm(0.3)
    start_x = (SLIDE_W - (box_w * 3 + box_gap * 2)) / 2
    boxes = [
        (SUCCESS, "🏆 Simple = 트래픽 수확기",
         "CTA 클릭률 55.4% 1위\n넓은 그물로 많이 끌어오는 메인"),
        (ACCENT, "🏆 Social = 질적 전환기",
         "Lead 전환율 9.6% 1위\n진짜 관심 있는 사람이 끝까지"),
        (WARN, "⚠ Fear = 중단 예정",
         "Lead 전환율 3.6% 꼴찌\n클릭은 되나 전환 안 됨"),
    ]
    for i, (color, t, b) in enumerate(boxes):
        bx = start_x + (box_w + box_gap) * i
        add_rect(s, bx, box_y, box_w, box_h, fill=SOFT, line=BORDER)
        add_rect(s, bx, box_y, Cm(0.3), box_h, fill=color)
        add_text(s, bx + Cm(0.5), box_y + Cm(0.2), box_w - Cm(0.7), Cm(0.9),
                 t, size=12, bold=True, color=NAVY)
        add_text(s, bx + Cm(0.5), box_y + Cm(1.1), box_w - Cm(0.7), box_h - Cm(1.3),
                 b, size=10, color=TEXT)

    # 연결 설명
    add_text(s, Cm(1.5), Cm(12.5), Cm(30.5), Cm(2.5),
             "결론: 3개 광고가 각기 다른 역할을 수행 · 버튼 많이 누르는 광고와 끝까지 전환되는 광고가 따로",
             size=12, color=TEXT)
    footer(s, 5)

    # ─── 6. Meta vs 데이터 분석 홍보 ───
    s = blank(prs)
    add_text(s, Cm(1.5), Cm(1), Cm(30), Cm(1.2),
             "Meta 광고 캠페인 비교 — AI-MBTI vs 데이터 분석 홍보",
             size=20, bold=True, color=NAVY)
    add_text(s, Cm(1.5), Cm(2.5), Cm(30.5), Cm(0.8),
             "같은 메타코드M 광고 계정 · 동일 기간(4/22~23) · 동일 조건(OUTCOME_LEADS)",
             size=11, color=MUTED)
    add_table(s, Cm(1.5), Cm(3.5), Cm(30.5), Cm(9),
              ["지표", "Simple (1위)", "Social", "Fear", "데이터 분석 홍보"],
              [
                  ["광고 목표", "잠재고객", "잠재고객", "잠재고객", "잠재고객"],
                  ["광고비", "₩29,891", "₩30,260", "₩33,593", "₩95,408"],
                  ["도달 (실사용자)", "1,526명", "1,242명", "1,495명", "2,627명"],
                  ["클릭률 (CTR)", "3.02%", "3.34%", "2.80%", "2.22%"],
                  ["클릭당 비용 (CPC)", "₩482", "₩593", "₩622", "₩1,403"],
                  ["Lead", "9명", "8명", "5명", "11명"],
                  ["1명 모시는 비용 (CPL)", "₩3,321 🏆", "₩3,783", "₩6,719", "₩8,673"],
              ],
              highlight=6)
    add_text(s, Cm(1.5), Cm(13), Cm(30.5), Cm(3),
             "• CPL: Simple ₩3,321 vs 데이터 분석 홍보 ₩8,673 → 2.6배 저렴\n"
             "• CPC: Simple ₩482 vs 데이터 분석 홍보 ₩1,403 → 2.9배 저렴\n"
             "• 효율 역전: 데이터 분석 홍보가 3배 더 쓰고도 Lead 11명 (Simple 9명과 거의 동일)",
             size=11, color=TEXT)
    footer(s, 6)

    # ─── 7. 대시보드 (데이터 분석 역량) ───
    s = blank(prs)
    add_text(s, Cm(1.5), Cm(1), Cm(30), Cm(1.2),
             "데이터 분석 역량 — 실제 운영 대시보드",
             size=20, bold=True, color=NAVY)
    add_image(s, "screen/dashboard3.png", x=Cm(1.5), y=Cm(2.8), w=Cm(30.5))
    footer(s, 7)

    # ─── 8. 서비스 흐름 ───
    s = blank(prs)
    add_text(s, Cm(1.5), Cm(1), Cm(30), Cm(1.2),
             "서비스 흐름 — 진단 → 결과 → 학습 연결",
             size=20, bold=True, color=NAVY)
    img_w = Cm(9.5)
    gap = Cm(1)
    start_x = (SLIDE_W - (img_w * 3 + gap * 2)) / 2
    for i, (rel, cap) in enumerate([
        ("screen/service-test.png", "1단계 · 20문항 진단"),
        ("screen/service-result.png", "2단계 · 16유형 결과"),
        ("screen/service-cta.png", "3단계 · 전자책·단톡방"),
    ]):
        add_image(s, rel, x=start_x + (img_w + gap) * i, y=Cm(3), w=img_w)
        add_text(s, start_x + (img_w + gap) * i, Cm(15.5), img_w, Cm(1),
                 cap, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 8)

    # ─── 9. Claude AI 평가 ───
    s = blank(prs)
    add_text(s, Cm(1.5), Cm(1), Cm(30), Cm(1.2),
             "Claude AI의 성과 평가",
             size=20, bold=True, color=NAVY)
    add_text(s, Cm(1.5), Cm(2.5), Cm(30.5), Cm(0.8),
             "AI(Claude)가 이 자료의 숫자를 따로 확인하고 매긴 종합 평가",
             size=11, color=MUTED)
    add_table(s, Cm(1.5), Cm(3.5), Cm(30.5), Cm(11),
              ["항목", "점수", "평가", "이유"],
              [
                  ["광고비 효율", "A+", "아주 잘함", "업계 평균의 1/6 비용, 클릭률 1.6배"],
                  ["광고 다양성", "A", "잘함", "3가지 광고 동시 테스트 비교"],
                  ["숫자 정확도", "A", "잘함", "페이스북 광고 관리자 숫자와 완전 일치"],
                  ["의사결정 속도", "B+", "괜찮음", "한 화면 대시보드로 빠르게 판단"],
                  ["전체 규모", "B", "이제 시작", "2일에 17명 — 예산 확장 검증 필요"],
                  ["종합", "A-", "광고 효율 입증", "규모 키우고 구매율 높이는 게 과제"],
              ],
              highlight=5)
    footer(s, 9)

    # ─── 10. 연락처 ───
    s = blank(prs)
    add_text(s, Cm(0), Cm(6), SLIDE_W, Cm(2),
             "연락처",
             size=40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Cm(0), Cm(9.5), SLIDE_W, Cm(1),
             "서비스: https://mcodegc.com",
             size=18, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Cm(0), Cm(11), SLIDE_W, Cm(1),
             "IR / PR 문의: information@mcode.co.kr",
             size=18, color=TEXT, align=PP_ALIGN.CENTER)
    footer(s, 10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"✓ 생성 완료: {OUT}")
    print(f"  크기: {OUT.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    build()
